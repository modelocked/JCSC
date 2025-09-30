#!/usr/bin/env python3
# -*- coding: utf-8 -*-
r"""
powerpoint_template_python.py

Build an Arnhem-styled deck from JSON using your template, with:
- Core fixed slides (title, declaration, contents, bottom line, conclusion title,
  conclusion, major references)
- Repeating pair per section: Section Title + Section Content (2-column)
- Forces white background; flips white text to black; locks font sizes (no autosize)

Usage:
    python powerpoint_template_python.py input.json Arnhem_Output.pptx
"""

import sys, os, json
from typing import List, Optional, Dict, Any
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.dml.color import RGBColor

TEMPLATE = r"C:\Users\User\Documents\CLEM\Arnhem_slide_template.pptx"

# Preferred layout names (rename your layouts once to these for exact matching).
LAYOUT_BY_PURPOSE = {
    "presentation_title": ["Arnhem – Presentation Title", "Title Slide", "Title", "Cover", "Title page"],
    "declaration":        ["Arnhem – Declaration", "Declaration", "Title and Content"],
    "contents":           ["Arnhem – Contents", "Contents", "Agenda", "Title and Content"],
    "bottom_line":        ["Arnhem – Bottom Line", "Bottom Line", "BLUF", "Title and Content"],
    "conclusion_title":   ["Arnhem – Conclusion Title", "Section Header", "Title Only"],
    "conclusion":         ["Arnhem – Conclusion", "Conclusion", "Title and Content"],
    "major_references":   ["Arnhem – Major References", "References", "Title and Content"],
    "section_title":      ["Arnhem – Section Title", "Section Header", "Title Only"],
    "section_content":    ["Arnhem – Section Content (2-col)", "Two Content", "Comparison", "Title and Content"]
}

# Font sizes you want (locked; no autosize)
S = {
    "title": Pt(36),
    "subtitle": Pt(20),
    "h1": Pt(28),
    "body": Pt(18),
    "sub1": Pt(16),
    "sub2": Pt(14),
}

WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)

# ------------------------ helpers ------------------------
def _layout_has(layout, t):
    for shp in layout.placeholders:
        try:
            if shp.placeholder_format.type == t:
                return True
        except Exception:
            pass
    return False

def _pick_layout(prs: Presentation, purpose: str):
    wanted = [n.strip().lower() for n in LAYOUT_BY_PURPOSE.get(purpose, [])]
    for lay in prs.slide_layouts:
        if lay.name and lay.name.strip().lower() in wanted:
            return lay
    # heuristics
    if purpose in ("presentation_title", "section_title", "conclusion_title"):
        for lay in prs.slide_layouts:
            if _layout_has(lay, PP_PLACEHOLDER.TITLE) and _layout_has(lay, PP_PLACEHOLDER.SUBTITLE):
                return lay
        for lay in prs.slide_layouts:
            if _layout_has(lay, PP_PLACEHOLDER.TITLE):
                return lay
    if purpose in ("declaration","contents","bottom_line","conclusion","major_references","section_content"):
        # prefer title+body or two content
        for lay in prs.slide_layouts:
            if _layout_has(lay, PP_PLACEHOLDER.TITLE) and _layout_has(lay, PP_PLACEHOLDER.BODY):
                return lay
        for lay in prs.slide_layouts:
            if _layout_has(lay, PP_PLACEHOLDER.BODY):
                return lay
    return prs.slide_layouts[0]

def _ph_by_type(slide, t):
    for shp in slide.placeholders:
        try:
            if shp.placeholder_format.type == t:
                return shp
        except Exception:
            pass
    return None

def _first_text_ph(slide, exclude_ids=set()):
    for shp in slide.placeholders:
        if id(shp) in exclude_ids:
            continue
        if hasattr(shp, "text_frame"):
            return shp
    return None

def _force_white_background(slide):
    # background fill
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE
    # also recolor any full-bleed rectangles (common in branded themes)
    sw, sh = slide.part.presentation.slide_width, slide.part.presentation.slide_height
    for shp in slide.shapes:
        try:
            if shp.width > sw * 0.95 and shp.height > sh * 0.95 and shp.fill:
                shp.fill.solid()
                shp.fill.fore_color.rgb = WHITE
        except Exception:
            pass

def _lock_no_autosize(tf):
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    tf.word_wrap = True

def _apply_text(tf, text: str, size: Pt, align=None, color=BLACK):
    _lock_no_autosize(tf)
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text or ""
    r.font.size = size
    r.font.color.rgb = BLACK if color == WHITE else color
    if align is not None:
        p.alignment = align

def _apply_bullets(tf, bullets, sizes=(S["body"], S["sub1"], S["sub2"]), color=BLACK):
    _lock_no_autosize(tf)
    tf.clear()
    def add_line(text, level):
        p = tf.paragraphs[0] if (len(tf.paragraphs)==1 and tf.paragraphs[0].text=="") else tf.add_paragraph()
        p.text = str(text)
        p.level = level
        for r in p.runs:
            r.font.size = sizes[min(level, len(sizes)-1)]
            r.font.color.rgb = BLACK if color == WHITE else color
    def walk(items, lvl=0):
        for it in items:
            if isinstance(it,(list,tuple)):
                walk(it, lvl+1)
            else:
                add_line(it, lvl)
    walk(bullets, 0)

def _flip_white_text_to_black(slide):
    # make sure nothing ends up white on white
    for shp in slide.shapes:
        if hasattr(shp, "text_frame") and shp.text_frame:
            for p in shp.text_frame.paragraphs:
                for r in p.runs:
                    try:
                        if r.font.color and r.font.color.rgb == WHITE:
                            r.font.color.rgb = BLACK
                    except Exception:
                        pass

# ------------------------ builders ------------------------
def add_presentation_title(prs, title, subtitle):
    s = prs.slides.add_slide(_pick_layout(prs, "presentation_title"))
    _force_white_background(s)
    t = s.shapes.title or _ph_by_type(s, PP_PLACEHOLDER.TITLE)
    if t and hasattr(t,"text_frame"): _apply_text(t.text_frame, title, S["title"], PP_ALIGN.LEFT)
    sub = _ph_by_type(s, PP_PLACEHOLDER.SUBTITLE) or _first_text_ph(s, {id(t)} if t else set())
    if subtitle and sub and hasattr(sub,"text_frame"): _apply_text(sub.text_frame, subtitle, S["subtitle"], PP_ALIGN.LEFT)
    _flip_white_text_to_black(s)

def add_title_body(prs, purpose_key, title, bullets=None):
    s = prs.slides.add_slide(_pick_layout(prs, purpose_key))
    _force_white_background(s)
    t = s.shapes.title or _ph_by_type(s, PP_PLACEHOLDER.TITLE)
    if t and hasattr(t,"text_frame"): _apply_text(t.text_frame, title, S["h1"], PP_ALIGN.LEFT)
    body = _ph_by_type(s, PP_PLACEHOLDER.BODY) or _first_text_ph(s, {id(t)} if t else set())
    if body and hasattr(body,"text_frame"):
        if bullets: _apply_bullets(body.text_frame, bullets)
        else: _apply_text(body.text_frame, "", S["body"])
    _flip_white_text_to_black(s)

def add_section_title(prs, title):
    s = prs.slides.add_slide(_pick_layout(prs, "section_title"))
    _force_white_background(s)
    t = s.shapes.title or _ph_by_type(s, PP_PLACEHOLDER.TITLE)
    if t and hasattr(t,"text_frame"): _apply_text(t.text_frame, title, S["h1"], PP_ALIGN.LEFT)
    else:
        box = s.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(9), Inches(1))
        _apply_text(box.text_frame, title, S["h1"], PP_ALIGN.LEFT)
    _flip_white_text_to_black(s)

def add_section_content(prs, title, left, right=None):
    s = prs.slides.add_slide(_pick_layout(prs, "section_content"))
    _force_white_background(s)
    t = s.shapes.title or _ph_by_type(s, PP_PLACEHOLDER.TITLE)
    if t and hasattr(t,"text_frame"): _apply_text(t.text_frame, title, S["h1"], PP_ALIGN.LEFT)

    # find two BODY placeholders if present
    bodies = []
    for ph in s.placeholders:
        try:
            if ph.placeholder_format.type == PP_PLACEHOLDER.BODY:
                bodies.append(ph)
        except Exception:
            pass
    left_ph  = bodies[0] if bodies else _first_text_ph(s, {id(t)} if t else set())
    right_ph = bodies[1] if len(bodies) > 1 else None

    if left_ph and hasattr(left_ph,"text_frame"):  _apply_bullets(left_ph.text_frame, left or [])
    if right and right_ph and hasattr(right_ph,"text_frame"):
        _apply_bullets(right_ph.text_frame, right)
    elif right and not right_ph and left_ph and hasattr(left_ph,"text_frame"):
        tf = left_ph.text_frame
        tf.add_paragraph()
        _apply_bullets(tf, right)

    _flip_white_text_to_black(s)

# ------------------------ pipeline ------------------------
def build_from_json(prs: Presentation, spec: Dict[str, Any]):
    core = spec.get("core", {})
    if "presentation_title" in core:
        x = core["presentation_title"]; add_presentation_title(prs, x.get("title",""), x.get("subtitle"))
    if "declaration" in core:
        x = core["declaration"];        add_title_body(prs, "declaration", x.get("title",""), x.get("bullets"))
    if "contents" in core:
        x = core["contents"];           add_title_body(prs, "contents", x.get("title","Contents"), x.get("bullets"))
    if "bottom_line" in core:
        x = core["bottom_line"];        add_title_body(prs, "bottom_line", x.get("title","Bottom Line"), x.get("bullets"))
    if "conclusion_title" in core:
        x = core["conclusion_title"];   add_title_body(prs, "conclusion_title", x.get("title","Conclusions"))
    if "conclusion" in core:
        x = core["conclusion"];         add_title_body(prs, "conclusion", x.get("title","Conclusion"), x.get("bullets"))
    if "major_references" in core:
        x = core["major_references"];   add_title_body(prs, "major_references", x.get("title","Major References"), x.get("bullets"))

    for sec in spec.get("sections", []):
        add_section_title(prs, sec.get("title",""))
        c = sec.get("content", {})
        add_section_content(prs, c.get("title",""), c.get("left",[]) or [], c.get("right") or None)

def main():
    if len(sys.argv) < 3:
        print("Usage: python powerpoint_template_python.py input.json Arnhem_Output.pptx")
        sys.exit(1)
    in_path, out_path = sys.argv[1], sys.argv[2]
    if not os.path.isfile(TEMPLATE):
        raise FileNotFoundError(f"Template not found: {TEMPLATE}")
    spec = json.load(open(in_path, "r", encoding="utf-8"))
    prs = Presentation(TEMPLATE)
    build_from_json(prs, spec)
    prs.save(out_path)
    print("Saved:", os.path.abspath(out_path))

if __name__ == "__main__":
    main()
